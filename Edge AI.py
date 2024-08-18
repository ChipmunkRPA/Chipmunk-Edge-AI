import os
import subprocess
import gradio as gr
import ollama
from langchain_community.embeddings import OllamaEmbeddings
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ConversationBufferMemory
from langchain_community.llms import Ollama
from langchain_chroma import Chroma
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
import chardet
import platform
import psutil
import webbrowser


# Disable ChromaDB telemetry
os.environ['CHROMA_TELEMETRY_ENABLED'] = 'false'

# Global variables
llm = None
embeddings = None
vector_store = None
chain = None
memory = None
use_case_folder = "use_cases"
ollama_process = None  # To hold the subprocess for Ollama

# Ensure the base directory exists
os.makedirs(use_case_folder, exist_ok=True)

css = """
    :root {
    --primary-color: #3498db; /* Primary color */
    --secondary-color: #2ecc71;
    --background-color: #E0F8E0; /* Light green background color */
    --text-color: #333;
    --border-color: #e0e0e0;
    --heading-font-family: 'Arial', sans-serif; /* Arial font for headings */
    --body-font-family: 'Arial', sans-serif; /* Arial font for body text */
    --neon-grey: #a9a9a9; /* Neon grey color */
}

.center-image {
    display: flex;
    justify-content: center;  /* Center the image horizontally */
    margin: 20px 0;  /* Add vertical margin if needed */
}

.center-image img {
    width: 100px;  /* Set the image width */
    height: 100px;  /* Set the image height */
    object-fit: cover;  /* Ensure the image fits within the specified dimensions */
}

body {
    font-family: var(--body-font-family); /* Use Arial for body text */
    color: var(--text-color);
}

/* Heading styles */
h3, h4, h5 {
    font-family: var(--heading-font-family); /* Use Arial for headings */
    font-weight: bold; /* Optional: Make headings bold */
    color: var(--primary-color); /* Set the color for headings to primary color */
}

h1, h2, h6 {
    font-family: var(--heading-font-family); /* Use Arial for headings */
    font-weight: bold; /* Optional: Make headings bold */
    text-align: center;
    color: #36454F; /* Set the color for headings to primary color */
    display:block;
}

/* Gradio container styles */
.gradio-container {
    max-width: 1200px;
    margin: 0 auto;
    padding: 20px;
    border: 2px solid var(--neon-grey); /* Neon grey frame */
    box-shadow: 0 0 10px var(--neon-grey), 0 0 20px var(--neon-grey), 0 0 30px var(--neon-grey); /* Neon glow effect */
    border-radius: 8px;
    background-color: #F0F0F0; /* Set the container background to light green */
}

/* Tab styles */
.tab-nav {
    background-color: transparent; /* Make tab container background transparent */
    border-radius: 8px;
    overflow: hidden;
    display: flex;
    justify-content: center;
    width: 200;
}

.tab-nav > div {
    flex: 1;
    text-align: center;
    padding: 10px;
    box-sizing: border-box;
}

.tab-nav button {
    margin: 0 5px; /* Add horizontal spacing between tabs */
    padding: 12px 20px;
    background-color: transparent; /* Make tab button background transparent */
    border: 2px solid var(--primary-color); /* Primary color border for tabs */
    border-radius: 4px;
    color: var(--primary-color);
    font-weight: bold; /* Set tab button text to bold */
    font-family: var(--body-font-family); /* Use Arial for tab buttons */
    transition: background-color 0.3s, border-color 0.3s, color 0.3s;
}

.tab-nav button:hover,
.tab-nav button.selected {
    background-color: var(--primary-color);
    color: white;
    border-color: var(--primary-color); /* Maintain primary color border on hover/selected */
}

/* Input and button styles */
input[type="text"],
textarea {
    width: 100%;
    padding: 12px;
    border: 1px solid var(--border-color);
    border-radius: 6px;
    font-size: 16px;
    font-family: var(--body-font-family); /* Use Arial for inputs and textareas */
    transition: border-color 0.3s;
}

input[type="text"]:focus,
textarea:focus {
    border-color: var(--primary-color);
    outline: none;
}

button {
    background-color: #FFFDD0;
    color: #3498db;
    border: none;
    padding: 12px 20px;
    border-radius: 6px;
    font-size: 16px;
    cursor: pointer;
    font-family: var(--body-font-family); /* Use Arial for buttons */
    transition: background-color 0.3s;
}

button:hover {
    background-color: #2980b9;
}

/* Chatbot styles */
.chat-message {
    background-color: white;
    border-radius: 8px;
    padding: 12px;
    margin-bottom: 10px;
    box-shadow: 0 2px 5px rgba(0, 0, 0, 0.05);
    font-family: var(--body-font-family); /* Use Arial for chat messages */
}

.user-message {
    background-color: #e8f5fe;
}

.bot-message {
    background-color: #f0f0f0;
}

/* Slider styles */
.slider {
    -webkit-appearance: none;
    width: 100%;
    height: 8px;
    border-radius: 5px;
    background: #d3d3d3;
    outline: none;
    opacity: 0.7;
    transition: opacity 0.2s;
}

.slider:hover {
    opacity: 1;
}

.slider::-webkit-slider-thumb {
    -webkit-appearance: none;
    appearance: none;
    width: 20px;
    height: 20px;
    border-radius: 50%;
    background: var(--secondary-color);
    cursor: pointer;
}

.slider::-moz-range-thumb {
    width: 20px;
    height: 20px;
    border-radius: 50%;
    background: var(--secondary-color);
    cursor: pointer;
}

/* Hide footer */
footer {
    display: none !important;
}

/* Responsive design */
@media (max-width: 768px) {
    .gradio-container {
        padding: 10px;
    }

    .tab-nav button {
        padding: 10px 15px;
    }
}
"""


def check_ollama_installed():
    """Check if the Ollama CLI is installed."""
    try:
        # Attempt to run a command using the Ollama CLI
        subprocess.run(["ollama", "--version"], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        return True
    except FileNotFoundError:
        return False
    except subprocess.CalledProcessError as e:
        print(f"Error checking Ollama installation: {e}")
        return False


def install_ollama():
    """Provide instructions for installing the Ollama CLI."""
    system = platform.system()

    if system == "Windows":
        return '\u274C'+" Please download and install Ollama from the official website: [Ollama Windows Installer](https://ollama.com/download)."
    elif system == "Darwin":  # macOS
        return '\u274C'+" Please run the following command in your terminal to install Ollama on macOS:\n\n`brew install ollama`"
    elif system == "Linux":
        return '\u274C'+" Please run the following commands in your terminal to install Ollama on Linux:\n\n`curl -fsSL https://ollama.com/install.sh | bash`"
    else:
        return '\u274C'+" Your operating system is not supported for automated installation of Ollama."


def check_ollama_running():
    """Check if the Ollama server is running."""

    def is_process_running(process_name):
        """Check if there is any running process that contains the given name."""
        for proc in psutil.process_iter(['name']):
            try:
                # Check if the process name contains the given name string
                if process_name.lower() in proc.info['name'].lower():
                    return True
            except (psutil.NoSuchProcess, psutil.AccessDenied, psutil.ZombieProcess):
                pass
        return False

    process_name = "ollama"
    if is_process_running(process_name):
        return '\u2705'+" Ollama is running."
    else:
        return '\u274C'+" Ollama server is not running."


def start_ollama_server():
    """Start the Ollama server if it's not already running."""
    global ollama_process
    status = check_ollama_running()
    if "not running" in status:
        return '\u274C'+ " Ollama is not running. Please start Ollama on your machine and restart the program."
    else:
        return '\u2705'+" Ollama is running correctly."


def download_models():
    """Download Ollama models."""
    try:
        # Pull the gemma2:2b model
        subprocess.run(["ollama", "pull", "gemma2:2b"], check=True)

        # Pull the nomic-embed-text model
        subprocess.run(["ollama", "pull", "nomic-embed-text"], check=True)

        return '\u2705'+" Models downloaded successfully."
    except subprocess.CalledProcessError as e:
        return f"Error occurred during model download: {e}"


def initialize_ai():
    """Initialize Ollama AI components."""
    global llm, embeddings, memory

    # Initialize Ollama client
    ollama_client = ollama.Client()

    # Initialize Ollama LLM and embeddings
    llm = Ollama(model="gemma2:2b")
    embeddings = OllamaEmbeddings(model="nomic-embed-text")

    # Initialize conversation memory
    memory = ConversationBufferMemory(memory_key="chat_history", return_messages=True)


def set_use_case(use_case, top_k=5):
    """Initialize AI and set use case."""
    global chain, vector_store
    if llm is None or embeddings is None:
        initialize_ai()

    # Load the Chroma vector store from the selected folder
    use_case_path = os.path.join(use_case_folder, use_case)
    vector_store = Chroma(
        persist_directory=use_case_path,
        embedding_function=embeddings,
    )

    # Reset the memory when changing use cases
    memory.clear()
    retriever = vector_store.as_retriever(search_type="similarity", search_kwargs={"k": top_k})
    chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=retriever,
        memory=memory,
    )
    return '\u2705'+f" Use case '{use_case}' set and AI initialized."


def parse_text_file(file_path):
    try:
        with open(file_path, 'rb') as file:
            raw_data = file.read()
            encoding = chardet.detect(raw_data)['encoding']
            content = raw_data.decode(encoding)
        return content
    except Exception as e:
        return f"Error parsing text file {file_path}: {e}"


def parse_docx_file(file_path):
    try:
        doc = docx.Document(file_path)
        content = "\n".join([para.text for para in doc.paragraphs])
        return content
    except Exception as e:
        return f"Error parsing docx file {file_path}: {e}"


def parse_pdf_file(file_path):
    try:
        doc = fitz.open(file_path)
        content = ""
        for page in doc:
            content += page.get_text()
        return content
    except Exception as e:
        return f"Error parsing pdf file {file_path}: {e}"


def parse_excel_file(file_path):
    try:
        df = pd.read_excel(file_path)
        content = df.to_string(index=False)
        return content
    except Exception as e:
        return f"Error parsing excel file {file_path}: {e}"


def parse_csv_file(file_path):
    try:
        df = pd.read_csv(file_path)
        content = df.to_string(index=False)
        return content
    except Exception as e:
        return f"Error parsing csv file {file_path}: {e}"


def parse_pptx_file(file_path):
    try:
        prs = Presentation(file_path)
        content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        return content
    except Exception as e:
        return f"Error parsing pptx file {file_path}: {e}"


def create_use_case(name, uploaded_files, chunk_size=512):
    global use_case_folder
    if embeddings is None:
        initialize_ai()

    chunks = []
    error_messages = []

    # Parse uploaded files
    for file_path in uploaded_files:
        try:
            if file_path.endswith('.txt'):
                content = parse_text_file(file_path)
            elif file_path.endswith('.docx'):
                content = parse_docx_file(file_path)
            elif file_path.endswith('.pdf'):
                content = parse_pdf_file(file_path)
            elif file_path.endswith('.xlsx'):
                content = parse_excel_file(file_path)
            elif file_path.endswith('.csv'):
                content = parse_csv_file(file_path)
            elif file_path.endswith('.pptx'):
                content = parse_pptx_file(file_path)
            else:
                error_message = f"Unsupported file format: {file_path}"
                error_messages.append(error_message)
                continue

            if isinstance(content, str) and content.startswith("Error"):
                error_messages.append(content)
                continue

            file_chunks = [content[i:i + chunk_size] for i in range(0, len(content), chunk_size)]
            chunks.extend(file_chunks)

        except Exception as e:
            error_message = f"Error processing file {file_path}: {e}"
            error_messages.append(error_message)

    # Create a directory for the new use case
    use_case_path = os.path.join(use_case_folder, name)
    os.makedirs(use_case_path, exist_ok=True)

    # Create a new Chroma vector store in the directory
    try:
        use_case_vector_store = Chroma(
            persist_directory=use_case_path,
            embedding_function=embeddings,
        )

        # Add the context to the vector store
        use_case_vector_store.add_texts(chunks)
    except Exception as e:
        error_message = f"Error creating or updating vector store: {e}"
        error_messages.append(error_message)
        return "Failed to create use case due to vector store error.", ""

    if error_messages:
        return "\n".join(error_messages), ""

    return '\u2705'+f" Use case '{name}' created successfully.", ""


def get_use_cases():
    # List all folders in the use case directory
    return [name for name in os.listdir(use_case_folder) if os.path.isdir(os.path.join(use_case_folder, name))]


def refresh_use_cases():
    """Refresh the list of use cases."""
    return gr.update(choices=get_use_cases())


def chat(message, history, creativity):
    global chain, llm, memory
    if chain is None:
        return "Please initialize the AI and set a use case first."
    try:
        prompt = "Please strictly follow the guidance below to answer the user questions:" \
                 "1. Only answer a user's question based on the context. If no context can be found, please ask a follow-up question." \
                 "2. Understand the keywords and objectives of the user question first and then find the closest answer from the context." \
                 "3. Please respond in a formal business style. Make sure your tone is formal and cheerful." \
                 "4. Consider the chat history before you answer questions." \
                 "5. Make sure your response is very detailed and list all the content you can find from the context in bullet points." \
                 "6. If in the end, you couldn't answer a user question, say sorry and direct the user to raysang@chipmunkrpa.com for further assistance." \
                 "The user question for you is below:"
        # Ensure history is formatted as a list of tuples
        formatted_history = [(user_msg, ai_msg) for user_msg, ai_msg in history]

        # Adjust the temperature based on the creativity slider
        llm.temperature = creativity

        # Update the memory with the current conversation
        for user_msg, ai_msg in formatted_history:
            memory.chat_memory.add_user_message(user_msg)
            memory.chat_memory.add_ai_message(ai_msg)

        # Generate a response using the chain
        response = chain.invoke({"question": prompt + message})

        # Add the new message to history
        new_history = history + [(message, response['answer'])]

        # Return the updated history to the chat window
        return new_history, ""
    except Exception as e:
        error_message = f"Error occurred during chat: {str(e)}"
        return [("Chipmunk Edge AI: ", error_message)], ""


def clear_history():
    global memory
    if memory is None:
        return "AI not initialized. Nothing to clear."
    memory.clear()
    return []


def export_history(history):
    with open("chat_history.txt", "w") as f:
        for human, ai in history:
            f.write(f"You: {human}\n\nEdge AI: {ai}\n\n")
    return "Chat history exported to chat_history.txt"


# Gradio interface
with gr.Blocks(css=css, title="Chipmunk Edge AI", theme=gr.themes.Soft()) as Edge_AI:
    gr.Markdown("# Chipmunk Edge AI 🤖 - The 100% Local AI System")
    gr.Markdown("## 🔥Powered by Chipmunk Robotics")

    with gr.Tab("Setup Models"):
        install_status = '\u2705'+" Ollama CLI is correctly installed." if check_ollama_installed() else install_ollama()
        server_status = start_ollama_server()
        gr.Markdown(f"**Ollama Installation Status:** {install_status}\n\n**Ollama Server Status:** {server_status}")
        with gr.Row():
            download_button = gr.Button("Download Models (One-time Only)")
            download_output = gr.Textbox(label="Download Status")
            download_button.click(download_models, outputs=download_output)

        gr.Markdown(
            """
            ### Chipmunk Edge AI Instructions:
            1. **Check Ollama CLI Installation:** Chipmunk Edge AI automatically checks if the Ollama CLI is properly installed and provides instructions if it's missing.
            2. **Ensure Ollama Server is Running:** Chipmunk Edge AI will verify whether an Ollama server is running. If not, please start ollama on your machine and restart this program.
            3. **Download Models:** In the "Setup Models" tab and click "Download Models" to set up the models locally. You only need to perform this once.
            4. **Create Use Case:** In the "Create Chipmunk Edge AI Use Case" tab, upload your local knowledge documents (when we say 'upload' here, it is really to process the documents in the local server within your own machine. No actual upload is made outside of your machine). Specify a chunk size if needed, and click "Create Use Case". Only pdf, txt, docx, pptx, xlsx and csv file formats are supported. Note more chuck size will give AI better ability to retrieve more information to answer your questions. However, it may reduce the accuracy and the speed of the response.
            5. **Refresh Use Cases:** In the "Chipmunk Edge AI" tab, use "Refresh Use Cases" to update the list of available use cases.
            6. **Chat:** In the "Chipmunk Edge AI" tab, select your use case, click "Initialize AI and Set Use Case", then ask your question, and click "Submit Question" to send it. It may take some time for the local AI model to respond depending on your machine's processing power.
            """
        )

    with gr.Tab("Create Chipmunk Edge AI Use Case"):
        gr.Markdown("**Note:** Supported file formats are PDF, DOCX, PPTX, TXT, XLSX, and CSV.")
        use_case_name = gr.Textbox(label="Use Case Name")
        file_upload = gr.File(label="Upload Documents", file_count="multiple", type="filepath")
        with gr.Row():
            create_button = gr.Button("Create Use Case")
            chunk_size_input = gr.Number(
                label="Chunk Size (Adjust This For Better Retrieval Results)",
                value=512)  # Input for chunk size
        create_output = gr.Textbox(label="Creation Status")
        create_button.click(
            create_use_case,
            inputs=[use_case_name, file_upload, chunk_size_input],
            outputs=[create_output, use_case_name]
        )

        gr.Markdown(
            """
            ### Chipmunk Edge AI Use Case:
            1. **Use Case Provides Context:** You can power up your own local AI model with extra and specific context only you own.
            2. **Ultimate Safety:** Everything is done within your own machine. No internet is required. No leakage is possible as long as your local machine is safe.
            """
        )

    with gr.Tab("Chipmunk Edge AI"):
        with gr.Row():
            use_case_dropdown = gr.Dropdown(choices=get_use_cases(), label="Select Use Case")
            set_use_case_output = gr.Textbox(label="Use Case Status")
        with gr.Row():
            refresh_button = gr.Button("Refresh Use Cases")
            refresh_button.click(refresh_use_cases, outputs=use_case_dropdown)
            set_use_case_button = gr.Button("Initialize AI and Set Use Case")
            set_use_case_button.click(set_use_case, inputs=[use_case_dropdown], outputs=set_use_case_output)
        chatbot = gr.Chatbot()
        msg = gr.Textbox(label="Your Question")
        with gr.Row():
            submit_button = gr.Button("Submit Question")
            clear = gr.Button("Clear History")
            export = gr.Button("Export History")

        creativity = gr.Slider(minimum=0, maximum=1, value=0.5, label="Creativity")

        # Update the submit button click event
        submit_button.click(chat, inputs=[msg, chatbot, creativity], outputs=[chatbot, msg])
        clear.click(clear_history, outputs=chatbot)
        export.click(export_history, inputs=chatbot)

        gr.Markdown(
            """
            ### Chipmunk Edge AI is sponsored by Chipmunk Robotics
            Check out our free resources! [Click here to learn more](https://www.chipmunkrpa.com/free-resource)

            Enjoy Chipmunk Edge AI so far? [Buy us some coffee](https://buy.stripe.com/fZe3dUez98iHfVC6oq) and [have a chat](https://calendly.com/raysang-chipmunkrpa/30min)
            """
        )

        gr.Markdown(
            """
            ### Chipmunk Edge AI Benefits:
            1. **Zero Cost:** Absolutely no third-party cost in running Chipmunk Edge AI.
            2. **Ultimate Safety:** Everything is done within your own machine. No internet is even required.
            3. **Best Small Model Used:** The best small LLM model is used without requiring the use of expensive machines or GPUs. Powered by Google Gemma 2:2b, equivalent to OpenAI GPT3.5.
            """
        )

    gr.Markdown(
        """
            <div class="center-image">
                <a href="https://chipmunkrpa.com" target="_blank">
                    <img src="https://picsum.photos/100" width="100" height="100" alt="Chipmunk Logo">
                </a>
            </div>        
        """
    )
    gr.Markdown(
        """
        ###### Copyright © 2024 Chipmunk Robotics Limited Liability Co. All Rights Reserved.          
        """
    )

app = Edge_AI.launch(share=False, prevent_thread_lock=True)

# Check if app is a tuple and extract the URL
if isinstance(app, tuple):
    app, local_url, share_url = app
    url = local_url
else:
    # Extract URL using known attributes
    try:
        url = f"http://{app.server_name}:{app.server_port}/"
    except:
        # If direct attributes don't exist, use a default or alternate method
        url = "http://127.0.0.1:7860/"  # Default Gradio port

# open the URL
webbrowser.open(url)

# Keep the application running
Edge_AI.block_thread()
